import os
from django.http import HttpResponse
from django.shortcuts import render, redirect
from django.views.decorators.http import require_POST
from django.core.files.storage import FileSystemStorage
from django.conf import settings
from .forms import VideoGenerationForm
import importlib.util
import sys

# Create your views here.
from django.contrib.auth import login as auth_login
from django.contrib.auth import authenticate
from .forms import CustomUserCreationForm

def home_page(request):
    return render(request, 'home/index.html')
def quiz_page(request):
    return render(request, 'home/quiz.html')
def notes_page(request):
    return render(request, 'home/notes.html')
def setting_page(request):
    return render(request, 'home/setting.html')
def interview_page(request):
    return render(request, 'home/interview.html')
def video_page(request):
    return render(request, 'home/video.html')
def signup(request):
    if request.method == 'POST':
        print("chal agay")
        form = CustomUserCreationForm(request.POST)
        if form.is_valid():
            user = form.save()
            auth_login(request, user)  # Use renamed import
            return redirect('login')
    else:
        form = CustomUserCreationForm()
    
    return render(request, 'home/signup.html', {'form': form})


def login(request):  # Renamed from 'login'
    if request.method == 'POST':
        email = request.POST.get('email')
        password = request.POST.get('password')
        user = authenticate(request, username=email, password=password)
        
        if user is not None:
            auth_login(request, user)  # Use renamed import
            return redirect('dashboard')
        else:
            return render(request, 'home/login.html', {'error': 'Invalid credentials'})
    
    return render(request, 'home/login.html')

def home(request):
    return render(request, 'home/dashboard.html')

VIDEO_LEARNING_DIR = os.path.join(settings.MEDIA_ROOT, 'Video_Learning')
os.makedirs(VIDEO_LEARNING_DIR, exist_ok=True)

def video_page(request):
    form = VideoGenerationForm()
    past_videos = get_past_videos()
    return render(request, 'home/video.html', {'form': form, 'past_videos': past_videos})

def get_past_videos():
    videos = []
    for filename in os.listdir(VIDEO_LEARNING_DIR):
        if filename.endswith('.mp4'):
            video_name = filename[:-4].replace('_', ' ').title()
            videos.append({'name': video_name, 'filename': filename})
    return videos

@require_POST
def generate_video(request):
    form = VideoGenerationForm(request.POST, request.FILES, request=request)
    if form.is_valid():
        topic = form.cleaned_data['topic']
        uploaded_file = request.FILES.get('files') # Get a single file

        extracted_text = ""
        fs = FileSystemStorage()
        temp_dir = os.path.join(settings.MEDIA_ROOT, 'temp_uploads')
        os.makedirs(temp_dir, exist_ok=True)

        try:
            file = uploaded_file # Use the single uploaded file
            file_path = os.path.join(temp_dir, file.name)
            fs.save(file_path, file)
            file_extension = os.path.splitext(file.name)[1].lower()

            if file_extension == '.pdf':
                try:
                    import PyPDF2
                    with open(file_path, 'rb') as pdf_file:
                        pdf_reader = PyPDF2.PdfReader(pdf_file)
                        extracted_text = ""
                        for page_num in range(len(pdf_reader.pages)):
                            page = pdf_reader.pages[page_num]
                            extracted_text += page.extract_text() + "\n\n"
                except ImportError:
                    return HttpResponse("Error: PyPDF2 is not installed. Please install it: pip install PyPDF2")
                except Exception as e:
                    return HttpResponse(f"Error during PDF processing with PyPDF2: {e}")
            elif file_extension in ['.doc', '.docx']:
                try:
                    from docx import Document
                    document = Document(file_path)
                    for paragraph in document.paragraphs:
                        extracted_text += paragraph.text + "\n"
                    extracted_text += "\n"
                except ImportError:
                    return HttpResponse("Error: python-docx is not installed. Please install it: pip install python-docx")
            os.remove(file_path) # Remove temporary file
        except Exception as e:
            return HttpResponse(f"Error processing file: {e}")
        finally:
            # Clean up the temporary directory
            for filename in os.listdir(temp_dir):
                file_path = os.path.join(temp_dir, filename)
                try:
                    if os.path.isfile(file_path):
                        os.remove(file_path)
                except Exception as e:
                    print(f"Error deleting temporary file: {e}")

        try:
            video_filename = f'{topic.replace(" ", "_")}_video.mp4'
            transcript_filename = f'{topic.replace(" ", "_")}_transcript.txt'
            video_output_path = os.path.join(VIDEO_LEARNING_DIR, video_filename)
            transcript_output_path = os.path.join(VIDEO_LEARNING_DIR, transcript_filename)

            # Import the video_gen module using importlib
            import importlib
            module_name = 'home.video_gen'
            module = importlib.import_module(module_name)

            script = generate_script(extracted_text, topic)

            # Call the video generation script
            transcript_content = module.main_with_input(topic, script, video_output_path, transcript_output_path)

            # Write the transcript file explicitly with UTF-8 encoding
            try:
                with open(transcript_output_path, 'w', encoding='utf-8') as f:
                    f.write(transcript_content)
            except Exception as e:
                print(f"Error writing transcript file: {e}")

            return render(request, 'home/video.html', {'form': VideoGenerationForm(), 'generation_success': True, 'past_videos': get_past_videos()})

        except ImportError as e:
            return HttpResponse(f"Error during video generation (import): {e}")
        except Exception as e:
            return HttpResponse(f"Error during video generation: {e}")
    else:
        return render(request, 'home/video.html', {'form': form, 'errors': form.errors, 'past_videos': get_past_videos()})

import google.generativeai as genai

def generate_script(extracted_text, topic):
    api = settings.GEMINI_API_KEY
    genai.configure(api_key=api)
    model = genai.GenerativeModel("gemini-2.0-flash")
    script = model.generate_content(
            f"For the given topic: {topic} "
            f"Generate a video script based on the following text: {extracted_text}. You can also use content from outside the text to explain. The video script should be in 1st person narrator explaining what the concept is without the use of any other person in the explanation unless needed explicitly. JUST GENERATE THE SCRIPT. DO NOT GENERATE ANY INTRODUCTION, EXPLANATION etc.",
        )
    return str(script)

from django.http import FileResponse, Http404

VIDEO_LEARNING_DIR = os.path.join(settings.BASE_DIR, 'Video_Learning')

def video_view(request, video_filename):
    video_path_on_disk = os.path.join(VIDEO_LEARNING_DIR, video_filename)
    transcript_filename = video_filename.replace('_video.mp4', '_transcript.txt')
    transcript_path = os.path.join(VIDEO_LEARNING_DIR, transcript_filename)

    # Extract video name from filename
    video_name = video_filename[:-4].replace('_', ' ').title()

    transcript_content = ""
    try:
        with open(transcript_path, 'r', encoding='utf-8') as f:
            transcript_content = f.read()
    except FileNotFoundError:
        transcript_content = "Transcript not found."
    except Exception as e:
        transcript_content = f"Error reading transcript: {e}"
    
    return render(request, 'home/video_view.html', {
        'video_filename': video_filename,
        'transcript': transcript_content,
        'video_name': video_name
    })

def serve_video(request, video_filename):
    video_path = os.path.join(VIDEO_LEARNING_DIR, video_filename)
    
    try:
        response = FileResponse(open(video_path, 'rb'), content_type='video/mp4')
        response['Content-Disposition'] = f'inline; filename={video_filename}'
        return response
    except FileNotFoundError:
        raise Http404("Video not found")


# Notes Module

import os
from django.shortcuts import render, redirect
from django.conf import settings
from django.http import HttpResponse
from django.views.decorators.http import require_POST
import google.generativeai as genai
from pdfminer.high_level import extract_text as pdf_extract_text
from docx import Document as DocxDocument
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import Paragraph
from reportlab.lib.units import inch
import regex as re
import logging
import tempfile
# ... other imports

logger = logging.getLogger(__name__) #create a logger.
@require_POST
def generate_notes(request):
    topic_name = request.POST.get('subject-title')
    uploaded_file = request.FILES.get('file-upload')
    additional_notes = request.POST.get('additional-notes')
    website_link = request.POST.get('link-input')

    extracted_text = ""

    if uploaded_file:
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()
        try:
            logger.info(f"File uploaded: {uploaded_file.name}, Extension: {file_extension}")
            if file_extension == '.pdf':
                try:
                    # Corrected line: Access the underlying file object
                    uploaded_file.file.seek(0)
                    with tempfile.NamedTemporaryFile(delete=False) as temp_pdf:
                        for chunk in uploaded_file.chunks():
                            temp_pdf.write(chunk)
                        temp_pdf_path = temp_pdf.name
                    extracted_text = pdf_extract_text(temp_pdf_path)
                    os.unlink(temp_pdf_path) #delete the temp file.
                    logger.info("PDF extracted successfully.")
                except Exception as pdf_error:
                    logger.error(f"PDF extraction error: {pdf_error}", exc_info=True)
                    return HttpResponse(f"Error processing PDF file. Please ensure it is a valid, text-based PDF.")
            elif file_extension == '.docx':
                doc = DocxDocument(uploaded_file)
                for paragraph in doc.paragraphs:
                    extracted_text += paragraph.text + "\n"
            else:
                return HttpResponse("Unsupported file type.")

        except Exception as e:
            logger.error(f"General error: {e}", exc_info=True)
            return HttpResponse(f"Error processing file: {e}")
    all_inputs = f"Topic: {topic_name}\n\n"
    if extracted_text:
        all_inputs += f"Syllabus Content:\n{extracted_text}\n\n"
    if additional_notes:
        all_inputs += f"Additional Notes:\n{additional_notes}\n\n"
    if website_link:
        all_inputs += f"External Website Link: {website_link}\n\n"

    api_key = settings.GEMINI_API_KEY
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel("gemini-2.0-flash-lite")

    try:
        prompt = f"""
        Generate comprehensive notes based on the following information:

        {all_inputs}

        Format the notes with:
        - Clear, large headings for main topics.
        - Numbered lists for points.
        - Clear paragraphs.
        - Remove all special characters, formatting codes, and extra whitespace from the generated text.
        - Clean the extracted text to create human readable notes.
        """
        response = model.generate_content(prompt)
        generated_notes = response.text
        # Remove extra characters.
        generated_notes = re.sub(r'[^a-zA-Z0-9\n\s.,:;-]', '', generated_notes)

    except Exception as e:
        return HttpResponse(f"Error generating notes with Gemini: {e}")

    notes_folder = 'notes_folder'
    os.makedirs(notes_folder, exist_ok=True)
    pdf_filename = f"{topic_name.replace(' ', '_')}_notes.pdf"
    pdf_filepath = os.path.join(notes_folder, pdf_filename)

    c = canvas.Canvas(pdf_filepath, pagesize=letter)
    styles = getSampleStyleSheet()
    normal_style = styles['Normal']
    heading_style = ParagraphStyle('Heading', parent=styles['Heading1'], fontSize=16, leading=20)

    y_position = 750

    lines = generated_notes.splitlines()
    for line in lines:
        if len(line.split()) > 0 and all(word.istitle() for word in line.split()): #A better heading detection.
            p = Paragraph(line, heading_style)
        else:
            p = Paragraph(line, normal_style)
        p_width, p_height = p.wrapOn(c, letter[0] - 2 * inch, 0)
        if y_position - p_height < 50:
            c.showPage()
            y_position = 750
        p.drawOn(c, inch, y_position - p_height)
        y_position -= p_height + 6

    c.save()

    return HttpResponse(f"Notes for '{topic_name}' generated successfully and saved as <a href='/download_notes/{pdf_filename}'>'{pdf_filename}'</a>")

#rest of the code remains the same.
def notes_page(request):
    return render(request, 'home/notes.html')

def download_notes(request, filename):
    file_path = os.path.join('notes_folder', filename)
    if os.path.exists(file_path):
        with open(file_path, 'rb') as f:
            response = HttpResponse(f.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
            return response
    return HttpResponse("File not found")



# Quiz Module

import os
import json
import uuid
from django.shortcuts import render, redirect
from django.conf import settings
from django.http import HttpResponse
from django.views.decorators.http import require_POST
import google.generativeai as genai
from docx import Document as DocxDocument
from pptx import Presentation as PPTXPresentation
from PyPDF2 import PdfReader

def quiz_page(request):
    return render(request, 'home/quiz.html')

import json
import re

def extract_json_from_response(response_text):
    json_match = None

    # Try to find JSON enclosed in ```json and ```
    code_block_match = re.search(r"```json\s*([\s\S]*?)\s*```", response_text)
    if code_block_match:
        json_match = code_block_match.group(1).strip()

    # If not found, try to find JSON directly (assuming it starts with [ or {)
    if not json_match:
        direct_match = re.search(r"(\[[\s\S]*?\]|\{[\s\S]*?\})", response_text)
        if direct_match:
            json_match = direct_match.group(1).strip()

    if json_match:
        try:
            return json.loads(json_match)
        except json.JSONDecodeError as e:
            print(f"Error decoding extracted JSON: {e}")
            print(f"Extracted JSON string: {json_match}")
            return None
    else:
        print("No JSON found in the response text.")
        print(f"Response Text: {response_text}")
        return None

@require_POST
def start_quiz(request):
    topic_name = request.POST.get('quiz-title')
    uploaded_file = request.FILES.get('file-upload')
    quiz_length = int(request.POST.get('quiz-length'))

    extracted_text = ""
    if uploaded_file:
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()
        try:
            if file_extension == '.pdf':
                pdf_reader = PdfReader(uploaded_file)
                for page in pdf_reader.pages:
                    extracted_text += page.extract_text() + "\n"
            elif file_extension == '.docx':
                doc = DocxDocument(uploaded_file)
                for paragraph in doc.paragraphs:
                    extracted_text += paragraph.text + "\n"
            elif file_extension == '.ppt' or file_extension == '.pptx':
                prs = PPTXPresentation(uploaded_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    extracted_text += run.text + "\n"
        except Exception as e:
            return HttpResponse(f"Error processing file: {e}")

    prompt = f"Generate {quiz_length} multiple-choice questions about '{topic_name}'. "
    if extracted_text:
        prompt += f"Use the following material as context:\n\n{extracted_text}\n\n"
    prompt += "Each question should have four options (A, B, C, D) and clearly indicate the correct answer in the JSON format: [{'question': '...', 'options': {'A': '...', 'B': '...', 'C': '...', 'D': '...'}, 'correct_answer': '...'}]"

    api_key = settings.GEMINI_API_KEY
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel("gemini-2.0-flash-lite")

    try:
        response = model.generate_content(prompt)
        print("Gemini API Response:")
        print(response.text)
        quiz_data = extract_json_from_response(response.text)

        if quiz_data:
            request.session['quiz_questions'] = quiz_data
            # Save quiz data temporarily
            quiz_session_id = uuid.uuid4()
            video_learning_folder = 'Video_Learning'
            os.makedirs(video_learning_folder, exist_ok=True)
            temp_quiz_file_path = os.path.join(video_learning_folder, f"quiz_{quiz_session_id}.json")
            with open(temp_quiz_file_path, 'w') as f:
                json.dump(quiz_data, f, indent=4)
            return redirect('quiz_attempt', quiz_session_id=quiz_session_id)
        else:
            return HttpResponse("Error: Could not extract valid JSON from the Gemini API response.")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")
        return HttpResponse(f"Error generating quiz: An unexpected error occurred. Check your server logs for details.")

def quiz_attempt(request, quiz_session_id):
    questions = request.session.get('quiz_questions')
    if not questions:
        return HttpResponse("Quiz questions not found.")
    context = {'questions': questions, 'quiz_session_id': quiz_session_id}
    return render(request, 'home/quizAttempt.html', context)

@require_POST
def submit_quiz(request):
    quiz_session_id = request.POST.get('quiz_session_id')
    questions = request.session.get('quiz_questions')
    if not questions:
        return HttpResponse("Quiz questions not found for submission.")

    score = 0
    total_questions = len(questions)
    user_answers = request.POST.dict()
    correct_answers = {}

    for i, question_data in enumerate(questions):
        correct_answers[f'question_{i+1}'] = question_data['correct_answer']
        user_answer = user_answers.get(f'question_{i+1}')
        if user_answer == question_data['correct_answer']:
            score += 1

    request.session['user_answers'] = user_answers
    request.session['correct_answers'] = correct_answers
    request.session['quiz_questions_display'] = questions # Keep questions for display

    video_learning_folder = 'Video_Learning'
    temp_quiz_file_path = os.path.join(video_learning_folder, f"quiz_{quiz_session_id}.json")
    try:
        os.remove(temp_quiz_file_path)
    except FileNotFoundError:
        print(f"Temporary quiz file not found: {temp_quiz_file_path}")
    except Exception as e:
        print(f"Error removing temporary quiz file: {e}")

    return redirect('quiz_result', score=score, total=total_questions)

def quiz_result(request, score, total):
    questions = request.session.get('quiz_questions_display', [])
    user_answers = request.session.get('user_answers', {})  # Ensure it's a dictionary
    formatted_answers = {f"question_{idx + 1}": str(answer) for idx, answer in enumerate(user_answers.values())}

    context = {
        'score': score,
        'total': total,
        'questions': questions,
        'user_answers': formatted_answers,
        'total_range': range(total)  # Add this to avoid filter issues
    }
    return render(request, 'home/quizResult.html', context)

